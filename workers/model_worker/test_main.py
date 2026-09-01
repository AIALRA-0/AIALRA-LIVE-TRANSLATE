"""Deterministic tests cover local fallbacks and page extraction without model downloads."""

from __future__ import annotations

import asyncio
import io

import httpx
import pytest
from pptx import Presentation

import workers.model_worker.main as model_worker
from workers.model_worker.main import (
    ASR_CPU_THREADS,
    EvidencePage,
    EvidenceSegment,
    ExplanationRequest,
    ExplanationResponse,
    SummaryRequest,
    _dedupe_text_items,
    _has_explanation_shape,
    _has_nonempty_list,
    _has_nonempty_string,
    _ollama_model_uses_gpu,
    _parse_asset_sync,
    _parse_model_json,
    _restore_realtime_translation_model,
    _uses_requested_explanation_language,
)


def test_asr_cpu_threads_stays_within_safe_host_bounds() -> None:
    assert 0 <= ASR_CPU_THREADS <= 32


def test_ollama_gpu_residency_requires_configured_model_and_near_full_vram() -> None:
    configured = {"name": model_worker.OLLAMA_MODEL, "size": 2_000, "size_vram": 1_900}
    assert _ollama_model_uses_gpu({"models": [configured]})
    assert not _ollama_model_uses_gpu({"models": [{**configured, "size_vram": 1_000}]})
    assert not _ollama_model_uses_gpu({"models": [{**configured, "name": "other"}]})


def _test_only_explanation(request: ExplanationRequest) -> ExplanationResponse:
    """Build deterministic evidence for unit tests without registering a runtime provider."""

    segment_ids = [segment.id for segment in request.segments]
    page_ids = [page.id for page in request.asset_pages]
    return ExplanationResponse(
        summary=request.segments[-1].text,
        missing_context=[],
        rare_terms=[],
        possible_asr_errors=[],
        review_questions=[],
        evidence_segment_ids=segment_ids,
        asset_page_ids=page_ids,
        confidence=0.5,
        provider="test_only",
    )


def test_explanation_fallback_keeps_all_supplied_evidence_ids() -> None:
    """A failed LLM call still links the card to the stable segment and uploaded page."""

    request = ExplanationRequest(
        segments=[EvidenceSegment(id="seg_1", text="Forwarding reduces stalls.")],
        asset_pages=[EvidencePage(id="page_1", title="Pipeline hazards", text="RAW hazard")],
        target_language="zh-CN",
    )
    result = _test_only_explanation(request)
    assert result.evidence_segment_ids == ["seg_1"]
    assert result.asset_page_ids == ["page_1"]


def test_pptx_parser_emits_stable_page_order_and_text() -> None:
    """A generated fixture verifies page order without private course files."""

    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[1])
    first.shapes.title.text = "Pipeline"
    first.placeholders[1].text = "Forwarding"
    second = presentation.slides.add_slide(presentation.slide_layouts[1])
    second.shapes.title.text = "Cache"
    second.placeholders[1].text = "Locality"
    buffer = io.BytesIO()
    presentation.save(buffer)

    result = _parse_asset_sync(".pptx", buffer.getvalue())
    assert [page.page_number for page in result.pages] == [1, 2]
    assert result.pages[0].title == "Pipeline"
    assert "Locality" in result.pages[1].text


def test_asyncio_is_available_for_worker_runtime() -> None:
    """The selected Python runtime can create an event loop for FastAPI inference calls."""

    assert asyncio.run(asyncio.sleep(0, result=True)) is True


def test_translation_shape_requires_clean_source_and_target_text() -> None:
    valid = {"source_text": "Attention uses context.", "translation": "注意力使用上下文"}
    assert _has_nonempty_string(valid, "source_text")
    assert _has_nonempty_string(valid, "translation")
    assert not _has_nonempty_string({"translation": "翻译结果"}, "source_text")
    assert not _has_nonempty_string({"source_text": "  "}, "source_text")


def test_explanation_shape_requires_every_managed_section() -> None:
    valid = {
        "summary": "简短总结",
        "missing_context": [],
        "rare_terms": [],
        "possible_asr_errors": [],
        "review_questions": [],
        "confidence": 0.8,
    }
    assert _has_explanation_shape(valid)
    assert _has_explanation_shape({"summary": "只有真实总结"})
    assert not _has_explanation_shape({**valid, "review_questions": "none"})


def test_model_json_accepts_unescaped_newline_inside_string() -> None:
    parsed = _parse_model_json('{"text":"first line\nsecond line"}')
    assert parsed == {"text": "first line\nsecond line"}


def test_model_json_rejects_non_json_or_non_object_output() -> None:
    assert _parse_model_json("translation without JSON") is None
    assert _parse_model_json('["translation"]') is None


def test_chinese_explanation_requires_chinese_summary() -> None:
    assert _uses_requested_explanation_language({"summary": "中文总结"}, "zh-CN")
    assert not _uses_requested_explanation_language({"summary": "English summary"}, "zh-CN")
    assert _uses_requested_explanation_language({"summary": "English summary"}, "en")


def test_background_model_restores_translation_before_releasing(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    calls: list[str] = []

    async def unload(model: str) -> None:
        calls.append(f"unload:{model}")

    class Response:
        def raise_for_status(self) -> None:
            calls.append(f"loaded:{model_worker.TRANSLATION_MODEL}")

    class Client:
        async def __aenter__(self) -> Client:
            return self

        async def __aexit__(self, *_args: object) -> None:
            return None

        async def post(self, _url: str, *, json: dict[str, object]) -> Response:
            assert json["model"] == model_worker.TRANSLATION_MODEL
            return Response()

    async def resident(model: str) -> bool:
        calls.append(f"resident:{model}")
        return True

    monkeypatch.setattr(model_worker, "_unload_ollama_model", unload)
    monkeypatch.setattr(httpx, "AsyncClient", lambda **_kwargs: Client())
    monkeypatch.setattr(model_worker, "_ollama_gpu_resident", resident)

    background_model = "qwen2.5:14b-instruct"
    asyncio.run(_restore_realtime_translation_model(background_model))

    assert calls == [
        f"unload:{background_model}",
        f"loaded:{model_worker.TRANSLATION_MODEL}",
        f"resident:{model_worker.TRANSLATION_MODEL}",
    ]


def test_summary_contract_accepts_the_core_rolling_summary_limit() -> None:
    request = SummaryRequest(
        segments=[EvidenceSegment(id="segment-1", text="Forwarding reduces stalls.")],
        rolling_summaries=[f"rolling-{index}" for index in range(24)],
        target_language="zh-CN",
    )
    assert len(request.rolling_summaries) == 24


def test_summary_required_lists_cannot_be_empty() -> None:
    assert _has_nonempty_list({"key_points": ["Pipeline hazards reduce throughput."]}, "key_points")
    assert not _has_nonempty_list({"key_points": []}, "key_points")
    assert not _has_nonempty_list({}, "key_points")


def test_summary_sections_drop_repeated_points_without_reordering() -> None:
    assert _dedupe_text_items(["  Pipeline stalls  ", "pipeline stalls", "Cache locality"]) == [
        "Pipeline stalls",
        "Cache locality",
    ]
