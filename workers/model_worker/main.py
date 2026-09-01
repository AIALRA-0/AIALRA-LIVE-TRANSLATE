"""Loopback model worker used by the Rust core through a versioned HTTP contract."""

from __future__ import annotations

import asyncio
import base64
import gc
import io
import json
import os
import sys
import threading
from collections.abc import Callable
from pathlib import Path
from typing import Annotated, Any

import httpx
import numpy as np
import numpy.typing as npt
from docx import Document
from fastapi import FastAPI, File, HTTPException, UploadFile
from PIL import Image
from pptx import Presentation
from pydantic import BaseModel, Field
from pypdf import PdfReader

_windows_dll_handles: list[Any] = []


def _register_bundled_cuda_dlls() -> None:
    """Expose uv-installed NVIDIA runtime DLLs to CTranslate2 on Windows."""

    if os.name != "nt" or not hasattr(os, "add_dll_directory"):
        return
    binary_directories: list[str] = []
    for entry in sys.path:
        package_root = Path(entry) / "nvidia"
        for package in ("cublas", "cudnn", "cuda_nvrtc"):
            binary_directory = package_root / package / "bin"
            if binary_directory.is_dir():
                resolved = str(binary_directory.resolve())
                if resolved not in binary_directories:
                    binary_directories.append(resolved)
                    _windows_dll_handles.append(os.add_dll_directory(resolved))
    if binary_directories:
        os.environ["PATH"] = os.pathsep.join(binary_directories + [os.environ.get("PATH", "")])


_register_bundled_cuda_dlls()

app = FastAPI(title="AIALRA Local Model Worker", version="1.0.0")

OLLAMA_URL = os.getenv("AIALRA_OLLAMA_URL", "http://127.0.0.1:11434").rstrip("/")
OLLAMA_MODEL = os.getenv("AIALRA_OLLAMA_MODEL", "qwen2.5:7b-instruct")
TRANSLATION_MODEL = os.getenv("AIALRA_TRANSLATION_MODEL", OLLAMA_MODEL)
EXPLANATION_MODEL = os.getenv("AIALRA_EXPLANATION_MODEL", "qwen2.5:7b-instruct")
SUMMARY_MODEL = os.getenv("AIALRA_SUMMARY_MODEL", "qwen2.5:14b-instruct")
# A 14B model can require a one-time cold CUDA load on a 16 GB card.  Keep that
# delay inside the explicit asynchronous summary lane instead of allowing the
# request to be cut off at the old 120-second transport timeout.  The core still
# keeps the session in ``processing`` until the real result arrives.
SUMMARY_TIMEOUT_SECONDS = max(
    30.0, min(float(os.getenv("AIALRA_SUMMARY_TIMEOUT_SECONDS", "120")), 120.0)
)
SUMMARY_CONTEXT_TOKENS = max(
    512, min(int(os.getenv("AIALRA_SUMMARY_CONTEXT_TOKENS", "2048")), 8192)
)
SUMMARY_MAX_TOKENS = max(
    160, min(int(os.getenv("AIALRA_SUMMARY_MAX_TOKENS", "320")), 600)
)
VISION_MODEL = os.getenv("AIALRA_VISION_MODEL", "qwen3-vl:8b-instruct")
ASR_MODEL_NAME = os.getenv("AIALRA_ASR_MODEL", "small")
ASR_DEVICE = os.getenv("AIALRA_ASR_DEVICE", "cuda")
ASR_COMPUTE_TYPE = os.getenv("AIALRA_ASR_COMPUTE_TYPE", "float16")
ASR_CPU_THREADS = max(0, min(32, int(os.getenv("AIALRA_ASR_CPU_THREADS", "12"))))
LLM_DEVICE = os.getenv("AIALRA_LLM_DEVICE", "cuda")

_asr_model: Any | None = None
_asr_lock = threading.Lock()


class HealthResponse(BaseModel):
    """Health separates worker availability from optional model readiness."""

    status: str
    asr_available: bool
    ollama_available: bool
    ollama_gpu_resident: bool
    model: str
    asr_provider: str
    asr_cpu_threads: int
    llm_provider: str


class AsrRequest(BaseModel):
    """PCM input avoids container ambiguity between browser, Android, and mini-app clients."""

    pcm_s16le_base64: str = Field(min_length=4, max_length=4_000_000)
    sample_rate: int = Field(ge=8_000, le=48_000)
    language: str = Field(min_length=2, max_length=32)
    initial_prompt: str = Field(default="", max_length=4_000)


class AsrResponse(BaseModel):
    """Final ASR result retains provider identity and measured audio duration."""

    text: str
    language: str
    confidence: float = Field(ge=0, le=1)
    duration_ms: int = Field(ge=0)
    provider: str


class GlossaryConstraint(BaseModel):
    """Confirmed terminology constrains translation without mutating source text."""

    source: str
    preferred: str
    do_not_translate: bool = False


class TranslationRequest(BaseModel):
    """Stable translation receives a bounded context and explicit terminology."""

    text: str = Field(min_length=1, max_length=20_000)
    source_language: str
    target_language: str
    glossary: list[GlossaryConstraint] = Field(default_factory=list, max_length=100)
    context: list[str] = Field(default_factory=list, max_length=10)


class TranslationResponse(BaseModel):
    """Provider identity lets the timeline verify the model and execution device."""

    source_text: str
    text: str
    provider: str


class EvidenceSegment(BaseModel):
    """Only stable segment IDs can become explanation evidence."""

    id: str
    text: str


class EvidencePage(BaseModel):
    """Parsed page text carries a stable page ID into the next explanation."""

    id: str
    title: str
    text: str


class ExplanationRequest(BaseModel):
    """The request contains a bounded transcript window and recent relevant pages."""

    segments: list[EvidenceSegment] = Field(min_length=1, max_length=20)
    asset_pages: list[EvidencePage] = Field(default_factory=list, max_length=12)
    target_language: str


class MissingContext(BaseModel):
    """Background additions retain the segment IDs that made them relevant."""

    text: str
    evidence_segment_ids: list[str]


class RareTerm(BaseModel):
    """A rare term receives one short explanation and traceable evidence."""

    term: str
    one_line: str
    evidence_segment_ids: list[str]
    asset_page_ids: list[str]


class ExplanationResponse(BaseModel):
    """Structured cards can be validated before entering the append-only timeline."""

    summary: str
    missing_context: list[MissingContext]
    rare_terms: list[RareTerm]
    possible_asr_errors: list[str]
    review_questions: list[str]
    evidence_segment_ids: list[str]
    asset_page_ids: list[str]
    confidence: float = Field(ge=0, le=1)
    provider: str


class SummaryRequest(BaseModel):
    """A final summary receives stable evidence only after the real-time queue drains."""

    segments: list[EvidenceSegment] = Field(min_length=1, max_length=64)
    asset_pages: list[EvidencePage] = Field(default_factory=list, max_length=24)
    rolling_summaries: list[str] = Field(default_factory=list, max_length=24)
    target_language: str


class SummaryResponse(BaseModel):
    """Every final summary remains traceable to stable transcript and page identifiers."""

    overview: str
    key_points: list[str]
    terminology: list[RareTerm]
    open_questions: list[str]
    evidence_segment_ids: list[str]
    asset_page_ids: list[str]
    provider: str


class ParsedPage(BaseModel):
    """Every page receives deterministic order, title, and extracted text."""

    page_number: int = Field(ge=1)
    title: str
    text: str


class AssetParseResponse(BaseModel):
    """Parser identity makes derived page text reproducible after upgrades."""

    parser: str
    pages: list[ParsedPage]


@app.get("/health", response_model=HealthResponse)
async def health() -> HealthResponse:
    """The core can keep recording when Ollama or faster-whisper is unavailable."""

    asr_available = _faster_whisper_importable()
    ollama_available = await _ollama_available()
    ollama_gpu_resident = await _ollama_gpu_resident()
    return HealthResponse(
        status="ok",
        asr_available=asr_available,
        ollama_available=ollama_available,
        ollama_gpu_resident=ollama_gpu_resident,
        model=OLLAMA_MODEL,
        asr_provider=f"faster-whisper:{ASR_MODEL_NAME}@{ASR_DEVICE}",
        asr_cpu_threads=ASR_CPU_THREADS,
        llm_provider=f"ollama:{OLLAMA_MODEL}@{LLM_DEVICE}",
    )


@app.post("/v1/asr/transcribe", response_model=AsrResponse)
async def transcribe(request: AsrRequest) -> AsrResponse:
    """ASR runs in a worker thread so model inference never blocks the HTTP event loop."""

    if not _faster_whisper_importable():
        raise HTTPException(status_code=503, detail="faster-whisper speech extra is unavailable")
    try:
        pcm_bytes = base64.b64decode(request.pcm_s16le_base64, validate=True)
    except ValueError as error:
        raise HTTPException(status_code=400, detail="invalid base64 PCM") from error
    if len(pcm_bytes) % 2:
        raise HTTPException(status_code=400, detail="PCM must contain complete 16-bit samples")
    audio = np.frombuffer(pcm_bytes, dtype="<i2").astype(np.float32) / 32768.0
    return await asyncio.to_thread(_transcribe_sync, audio, request)


@app.post("/v1/translate", response_model=TranslationResponse)
async def translate(request: TranslationRequest) -> TranslationResponse:
    """A translation is accepted only when the configured local Ollama model returns valid JSON."""

    glossary_lines = [
        f"{item.source} => {item.source if item.do_not_translate else item.preferred}"
        for item in request.glossary
    ]
    system = (
        "You clean and translate one provisional ASR lecture paragraph. Return JSON only. "
        "source_text must stay in the source language: restore punctuation and casing, join "
        "clearly split clauses, and remove only immediate accidental repetition from adjacent "
        "ASR windows. Never paraphrase, translate, invent, or remove meaning in source_text. "
        "translation must be one natural, meaning-based target-language paragraph that preserves "
        "the cleaned source's logical connections without adding facts. "
        "Context is previous-course context for terminology "
        "and pronoun resolution only: never translate, quote, or repeat Context. "
        "Preserve formulas, code, model numbers, and do-not-translate terms. "
        "Do not return labels or commentary outside JSON."
    )
    user = (
        f"Source language: {request.source_language}\n"
        f"Target language: {request.target_language}\n"
        f"Context: {' | '.join(request.context[-3:])}\n"
        f"Glossary: {'; '.join(glossary_lines)}\n"
        f"Text: {request.text}"
    )
    result = await _ollama_json(
        system,
        user,
        {
            "type": "object",
            "properties": {
                "source_text": {"type": "string", "maxLength": 20_000},
                "translation": {"type": "string", "maxLength": 20_000},
            },
            "required": ["source_text", "translation"],
            "additionalProperties": False,
        },
        max_tokens=1200,
        model=TRANSLATION_MODEL,
        timeout_seconds=45.0,
        attempts=2,
        accept=lambda payload: _has_nonempty_string(payload, "source_text")
        and _has_nonempty_string(payload, "translation"),
    )
    if isinstance(result, dict):
        return TranslationResponse(
            source_text=str(result["source_text"]).strip(),
            text=str(result["translation"]).strip(),
            provider=f"ollama:{TRANSLATION_MODEL}@{LLM_DEVICE}",
        )
    raise HTTPException(status_code=503, detail="local Ollama translation is unavailable")


@app.post("/v1/explain", response_model=ExplanationResponse)
async def explain(request: ExplanationRequest) -> ExplanationResponse:
    """The model writes bounded teaching content while trusted code attaches evidence IDs."""

    await _unload_ollama_model(VISION_MODEL)
    await _unload_ollama_model(SUMMARY_MODEL)
    segment_ids = [segment.id for segment in request.segments]
    page_ids = [page.id for page in request.asset_pages]
    system = (
        "You are a lecture comprehension assistant. Return compact JSON in the requested language. "
        "When target_language starts with zh, write every natural-language field "
        "in Simplified Chinese. "
        "Separate course statements from background knowledge. Do not repeat identifiers. "
        "Write a one-sentence summary. Return at most two missing-context items, three rare terms, "
        "two possible ASR errors, and two review questions. Explain each rare term in one sentence "
        "and keep every item concise."
    )
    language_instruction = (
        "All natural-language output fields must use Simplified Chinese.\n"
        if request.target_language.lower().startswith("zh")
        else ""
    )
    user = language_instruction + json.dumps(
        {
            "target_language": request.target_language,
            "segments": [segment.text for segment in request.segments],
            "asset_pages": [
                {"title": page.title, "text": page.text} for page in request.asset_pages
            ],
            "required_shape": {
                "summary": "string",
                "missing_context": ["string"],
                "rare_terms": [{"term": "string", "one_line": "string"}],
                "possible_asr_errors": ["string"],
                "review_questions": ["string"],
                "confidence": 0.0,
            },
        },
        ensure_ascii=False,
    )
    try:
        result = await _ollama_json(
            system,
            user,
            {
            "type": "object",
            "properties": {
                "summary": {"type": "string", "maxLength": 300},
                "missing_context": {
                    "type": "array",
                    "maxItems": 2,
                    "items": {"type": "string", "maxLength": 240},
                },
                "rare_terms": {
                    "type": "array",
                    "maxItems": 3,
                    "items": {
                        "type": "object",
                        "properties": {
                            "term": {"type": "string", "maxLength": 80},
                            "one_line": {"type": "string", "maxLength": 240},
                        },
                        "required": ["term", "one_line"],
                        "additionalProperties": False,
                    },
                },
                "possible_asr_errors": {
                    "type": "array",
                    "maxItems": 2,
                    "items": {"type": "string", "maxLength": 200},
                },
                "review_questions": {
                    "type": "array",
                    "maxItems": 2,
                    "items": {"type": "string", "maxLength": 200},
                },
                "confidence": {"type": "number", "minimum": 0, "maximum": 1},
            },
            "required": [
                "summary",
                "missing_context",
                "rare_terms",
                "possible_asr_errors",
                "review_questions",
                "confidence",
            ],
            "additionalProperties": False,
            },
            max_tokens=512,
            model=EXPLANATION_MODEL,
            timeout_seconds=75.0,
            attempts=2,
            accept=lambda payload: _has_explanation_shape(payload)
            and _uses_requested_explanation_language(payload, request.target_language),
        )
    finally:
        await _restore_realtime_translation_model(EXPLANATION_MODEL)
    if isinstance(result, dict):
        compact = {
            **result,
            "missing_context": [
                {"text": item, "evidence_segment_ids": segment_ids[-2:]}
                for item in result.get("missing_context", [])
                if isinstance(item, str)
            ],
            "rare_terms": [
                {
                    "term": item.get("term", ""),
                    "one_line": item.get("one_line", ""),
                    "evidence_segment_ids": segment_ids[-2:],
                    "asset_page_ids": page_ids,
                }
                for item in result.get("rare_terms", [])
                if isinstance(item, dict)
            ],
            "evidence_segment_ids": segment_ids,
            "asset_page_ids": page_ids,
        }
        normalized = _normalize_explanation(compact, segment_ids, page_ids)
        if normalized is not None:
            normalized.provider = f"ollama:{EXPLANATION_MODEL}@{LLM_DEVICE}"
            return normalized
    raise HTTPException(status_code=503, detail="local Ollama explanation is unavailable")


@app.post("/v1/summarize", response_model=SummaryResponse)
async def summarize(request: SummaryRequest) -> SummaryResponse:
    """The larger local model produces one evidence-bounded summary after recording stops."""

    await _unload_ollama_model(VISION_MODEL)
    await _unload_ollama_model(EXPLANATION_MODEL)
    if SUMMARY_MODEL != TRANSLATION_MODEL:
        # Remove the resident realtime model before loading 14B.  Relying on
        # Ollama's eviction heuristics made the one-shot path sensitive to the
        # exact lecture history and left too little headroom for CUDA ASR.
        await _unload_ollama_model(TRANSLATION_MODEL)
    # Whisper keeps a CUDA model resident during a recording.  Release it before
    # loading the one-shot summary model so the 16 GB card does not spend the
    # entire timeout evicting ASR allocations while Ollama is still cold-starting.
    await asyncio.to_thread(_release_asr_model_sync)
    segment_ids = [segment.id for segment in request.segments]
    page_ids = [page.id for page in request.asset_pages]
    system = (
        "You summarize a completed lecture using only supplied evidence. Return compact JSON. "
        "Do not invent facts, citations, or identifiers. Prefer the rolling summaries for the "
        "overall structure and use the sampled segments to verify details. Keep the overview "
        "under 300 characters, and return at most five key points, five terms, and three open "
        "questions."
    )
    if request.target_language.lower().startswith("zh"):
        system += " Write every natural-language field in Simplified Chinese."
    user = json.dumps(
        {
            "segments": [{"id": item.id, "text": item.text} for item in request.segments],
            "asset_pages": [
                {"id": item.id, "title": item.title, "text": item.text}
                for item in request.asset_pages
            ],
            "rolling_summaries": request.rolling_summaries,
        },
        ensure_ascii=False,
    )
    try:
        result = await _ollama_json(
            system,
            user,
            {
            "type": "object",
            "properties": {
                "overview": {"type": "string", "maxLength": 300},
                "key_points": {
                    "type": "array",
                    "minItems": 1,
                    "maxItems": 5,
                    "items": {"type": "string", "maxLength": 180},
                },
                "terminology": {
                    "type": "array",
                    "maxItems": 5,
                    "items": {
                        "type": "object",
                        "properties": {
                            "term": {"type": "string", "maxLength": 80},
                            "one_line": {"type": "string", "maxLength": 180},
                        },
                        "required": ["term", "one_line"],
                        "additionalProperties": False,
                    },
                },
                "open_questions": {
                    "type": "array",
                    "maxItems": 3,
                    "items": {"type": "string", "maxLength": 180},
                },
            },
            "required": ["overview", "key_points", "terminology", "open_questions"],
            "additionalProperties": False,
            },
            max_tokens=SUMMARY_MAX_TOKENS,
            model=SUMMARY_MODEL,
            unload_after=True,
            timeout_seconds=SUMMARY_TIMEOUT_SECONDS,
            num_ctx=SUMMARY_CONTEXT_TOKENS,
            attempts=1,
            accept=lambda payload: _has_nonempty_string(payload, "overview")
            and _has_nonempty_list(payload, "key_points"),
        )
    finally:
        await _restore_realtime_translation_model(SUMMARY_MODEL)
    if not isinstance(result, dict):
        raise HTTPException(status_code=503, detail="local Ollama summary is unavailable")
    terminology: list[RareTerm] = []
    seen_terms: set[str] = set()
    for item in result.get("terminology", []):
        if not isinstance(item, dict):
            continue
        term = str(item.get("term", "")).strip()
        one_line = str(item.get("one_line", "")).strip()
        term_key = " ".join(term.split()).casefold()
        if not term or not one_line or term_key in seen_terms:
            continue
        seen_terms.add(term_key)
        terminology.append(
            RareTerm(
                term=term,
                one_line=one_line,
                evidence_segment_ids=segment_ids,
                asset_page_ids=page_ids,
            )
        )
    return SummaryResponse(
        overview=str(result.get("overview", "")),
        key_points=_dedupe_text_items(result.get("key_points", [])),
        terminology=terminology,
        open_questions=_dedupe_text_items(result.get("open_questions", [])),
        evidence_segment_ids=segment_ids,
        asset_page_ids=page_ids,
        provider=f"ollama:{SUMMARY_MODEL}@{LLM_DEVICE}",
    )


@app.post("/v1/assets/parse", response_model=AssetParseResponse)
async def parse_asset(file: Annotated[UploadFile, File()]) -> AssetParseResponse:
    """Parsers receive in-memory bytes and never trust an uploaded path or archive member name."""

    data = await file.read()
    if len(data) > 50 * 1024 * 1024:
        raise HTTPException(status_code=413, detail="asset exceeds 50 MiB bootstrap limit")
    suffix = Path(file.filename or "asset.bin").suffix.lower()
    try:
        if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
            return await _parse_image_with_vlm(data)
        return await asyncio.to_thread(_parse_asset_sync, suffix, data)
    except (OSError, ValueError, KeyError) as error:
        raise HTTPException(status_code=422, detail="asset parser rejected the file") from error


def _faster_whisper_importable() -> bool:
    """Import probing avoids loading model weights during health checks."""

    try:
        import faster_whisper  # noqa: F401
    except ImportError:
        return False
    return True


def _get_asr_model() -> Any:
    """One lazy model instance preserves VRAM and serializes first-load races."""

    global _asr_model
    with _asr_lock:
        if _asr_model is None:
            from faster_whisper import WhisperModel

            _asr_model = WhisperModel(
                ASR_MODEL_NAME,
                device=ASR_DEVICE,
                compute_type=ASR_COMPUTE_TYPE,
                cpu_threads=ASR_CPU_THREADS,
            )
        return _asr_model


def _transcribe_sync(audio: npt.NDArray[np.float32], request: AsrRequest) -> AsrResponse:
    """A bounded four-second window produces a stable bootstrap segment."""

    model = _get_asr_model()
    language = None if request.language in {"auto", "mixed", "zh-en"} else request.language
    segments, info = model.transcribe(
        audio,
        language=language,
        initial_prompt=request.initial_prompt or None,
        beam_size=1,
        best_of=1,
        vad_filter=True,
        condition_on_previous_text=False,
    )
    realized = list(segments)
    text = " ".join(segment.text.strip() for segment in realized).strip()
    probabilities = [float(np.exp(segment.avg_logprob)) for segment in realized]
    confidence = float(np.clip(np.mean(probabilities), 0.0, 1.0)) if probabilities else 0.0
    duration_ms = int(len(audio) * 1_000 / request.sample_rate)
    return AsrResponse(
        text=text,
        language=str(getattr(info, "language", language or "unknown")),
        confidence=confidence,
        duration_ms=duration_ms,
        provider=f"faster-whisper:{ASR_MODEL_NAME}@{ASR_DEVICE}",
    )


async def _ollama_available() -> bool:
    """A short timeout prevents health checks from delaying the recording controls."""

    try:
        async with httpx.AsyncClient(timeout=1.0) as client:
            response = await client.get(f"{OLLAMA_URL}/api/tags")
            return response.is_success
    except httpx.HTTPError:
        return False


def _ollama_model_uses_gpu(payload: object, model: str = OLLAMA_MODEL) -> bool:
    """Accept the configured model only when at least 90 percent is resident in VRAM."""

    if not isinstance(payload, dict) or not isinstance(payload.get("models"), list):
        return False
    for item in payload["models"]:
        if not isinstance(item, dict) or item.get("name") != model:
            continue
        size = item.get("size")
        size_vram = item.get("size_vram")
        if (
            isinstance(size, int)
            and not isinstance(size, bool)
            and size > 0
            and isinstance(size_vram, int)
            and not isinstance(size_vram, bool)
        ):
            return size_vram >= size * 0.9
    return False


async def _ollama_gpu_resident(model: str = OLLAMA_MODEL) -> bool:
    """Read Ollama's process inventory instead of trusting a configured provider suffix."""

    try:
        async with httpx.AsyncClient(timeout=2.0) as client:
            response = await client.get(f"{OLLAMA_URL}/api/ps")
            response.raise_for_status()
            return _ollama_model_uses_gpu(response.json(), model)
    except (httpx.HTTPError, TypeError, ValueError):
        return False


async def _ollama_json(
    system: str,
    user: str,
    schema: dict[str, Any] | None = None,
    *,
    max_tokens: int = 768,
    accept: Callable[[dict[str, Any]], bool] | None = None,
    model: str = OLLAMA_MODEL,
    unload_after: bool = False,
    timeout_seconds: float = 90.0,
    attempts: int = 2,
    num_ctx: int | None = None,
) -> dict[str, Any] | None:
    """Ollama receives only text already allowed by the local session policy."""

    for attempt in range(max(1, min(attempts, 3))):
        try:
            async with httpx.AsyncClient(timeout=timeout_seconds) as client:
                repair = (
                    "\nThe prior response was invalid. Return exactly the requested JSON shape "
                    "with every required field populated."
                    if attempt
                    else ""
                )
                options: dict[str, Any] = {
                    "temperature": 0 if attempt == 0 else 0.1,
                    "seed": attempt,
                    "num_predict": max_tokens,
                }
                if num_ctx is not None:
                    options["num_ctx"] = num_ctx
                response = await client.post(
                    f"{OLLAMA_URL}/api/chat",
                    json={
                        "model": model,
                        "stream": False,
                        "format": schema or "json",
                        "messages": [
                            {"role": "system", "content": system},
                            {"role": "user", "content": f"{user}{repair}"},
                        ],
                        "options": options,
                    },
                )
                response.raise_for_status()
                content = response.json()["message"]["content"]
                parsed = _parse_model_json(content)
                if isinstance(parsed, dict) and (accept is None or accept(parsed)):
                    if unload_after:
                        await _unload_ollama_model(model)
                    return parsed
        except (httpx.HTTPError, KeyError, TypeError):
            await asyncio.sleep(1)
    if unload_after:
        await _unload_ollama_model(model)
    return None


async def _unload_ollama_model(model: str) -> None:
    """Release large one-shot models after proof-backed inference so ASR keeps VRAM headroom."""

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            response = await client.post(
                f"{OLLAMA_URL}/api/generate",
                json={"model": model, "keep_alive": 0, "stream": False},
            )
            response.raise_for_status()
    except httpx.HTTPError:
        return


async def _restore_realtime_translation_model(background_model: str) -> None:
    """Restore the low-latency model before a background lane releases its lock."""

    if background_model == TRANSLATION_MODEL:
        return
    await _unload_ollama_model(background_model)
    try:
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.post(
                f"{OLLAMA_URL}/api/generate",
                json={
                    "model": TRANSLATION_MODEL,
                    "prompt": "",
                    "keep_alive": -1,
                    "stream": False,
                },
            )
            response.raise_for_status()
        if not await _ollama_gpu_resident(TRANSLATION_MODEL):
            raise RuntimeError("translation model did not return to the GPU")
    except (httpx.HTTPError, RuntimeError) as error:
        raise HTTPException(
            status_code=503,
            detail="local translation model could not be restored after background inference",
        ) from error


def _release_asr_model_sync() -> None:
    """One-shot large models reclaim ASR VRAM only after the serialized queue is idle."""

    global _asr_model
    with _asr_lock:
        _asr_model = None
    gc.collect()


async def _ollama_text(
    system: str,
    user: str,
    *,
    max_tokens: int,
    model: str = OLLAMA_MODEL,
) -> str | None:
    """A plain local model response avoids fragile JSON quoting for translated prose."""

    for attempt in range(2):
        try:
            async with httpx.AsyncClient(timeout=45.0) as client:
                repair = (
                    "\nThe prior response was empty. Return only the translated text."
                    if attempt
                    else ""
                )
                response = await client.post(
                    f"{OLLAMA_URL}/api/chat",
                    json={
                        "model": model,
                        "stream": False,
                        "messages": [
                            {"role": "system", "content": system},
                            {"role": "user", "content": f"{user}{repair}"},
                        ],
                        "options": {
                            "temperature": 0 if attempt == 0 else 0.1,
                            "seed": attempt,
                            "num_predict": max_tokens,
                        },
                    },
                )
                response.raise_for_status()
                content = response.json()["message"]["content"]
                if isinstance(content, str) and content.strip():
                    return content.strip()
        except (httpx.HTTPError, KeyError, TypeError):
            await asyncio.sleep(1)
    return None


async def _parse_image_with_vlm(data: bytes) -> AssetParseResponse:
    """Local VLM performs OCR and visual explanation only when the image stays on this host."""

    with Image.open(io.BytesIO(data)) as image:
        image.verify()
    await _unload_ollama_model(EXPLANATION_MODEL)
    await _unload_ollama_model(SUMMARY_MODEL)
    if VISION_MODEL != TRANSLATION_MODEL:
        await _unload_ollama_model(TRANSLATION_MODEL)
    await asyncio.to_thread(_release_asr_model_sync)
    schema = {
        "type": "object",
        "properties": {
            "title": {"type": "string", "maxLength": 120},
            "ocr_text": {"type": "string", "maxLength": 8000},
            "visual_summary": {"type": "string", "maxLength": 1200},
            "key_terms": {
                "type": "array",
                "maxItems": 12,
                "items": {"type": "string", "maxLength": 100},
            },
        },
        "required": ["title", "ocr_text", "visual_summary", "key_terms"],
        "additionalProperties": False,
    }
    try:
        try:
            async with httpx.AsyncClient(timeout=240.0) as client:
                response = await client.post(
                    f"{OLLAMA_URL}/api/chat",
                    json={
                        "model": VISION_MODEL,
                        "stream": False,
                        "format": schema,
                        "messages": [
                            {
                                "role": "user",
                                "content": (
                                    "Extract all readable text and explain this lecture image in "
                                    "concise Simplified Chinese. Distinguish visible text from "
                                    "interpretation."
                                ),
                                "images": [base64.b64encode(data).decode("ascii")],
                            }
                        ],
                        "options": {"temperature": 0, "num_predict": 1200},
                    },
                )
                response.raise_for_status()
                result = _parse_model_json(str(response.json()["message"]["content"]))
        except (httpx.HTTPError, KeyError, TypeError, ValueError) as error:
            raise HTTPException(
                status_code=503,
                detail="local Ollama vision model is unavailable",
            ) from error
        if not isinstance(result, dict) or not _has_nonempty_string(result, "visual_summary"):
            raise HTTPException(status_code=503, detail="local Ollama vision result is invalid")
        gpu_resident = await _ollama_gpu_resident(VISION_MODEL)
        if not gpu_resident:
            raise HTTPException(
                status_code=503,
                detail="local Ollama vision model did not prove GPU residency",
            )
        title = str(result.get("title") or "课程图片")
        ocr_text = str(result.get("ocr_text") or "").strip()
        summary = str(result.get("visual_summary") or "").strip()
        terms = [
            str(item).strip()
            for item in result.get("key_terms", [])
            if isinstance(item, str) and item.strip()
        ]
        sections = [
            f"可见文字\n{ocr_text}" if ocr_text else "",
            f"图片解释\n{summary}",
            f"关键词\n{'、'.join(terms)}" if terms else "",
        ]
        return AssetParseResponse(
            parser=f"ollama:{VISION_MODEL}@{LLM_DEVICE}",
            pages=[
                ParsedPage(
                    page_number=1,
                    title=title,
                    text="\n\n".join(section for section in sections if section),
                )
            ],
        )
    finally:
        # Release VLM weights and restore the resident realtime model even when
        # inference, JSON validation, or the GPU residency proof fails.
        if VISION_MODEL != TRANSLATION_MODEL:
            await _unload_ollama_model(VISION_MODEL)
            await _restore_realtime_translation_model(VISION_MODEL)


def _parse_model_json(content: str) -> dict[str, Any] | None:
    """Accept unescaped control characters while retaining the JSON object boundary."""

    for strict in (True, False):
        try:
            parsed = json.loads(content, strict=strict)
        except json.JSONDecodeError:
            continue
        return parsed if isinstance(parsed, dict) else None
    return None


def _has_nonempty_string(payload: dict[str, Any], field: str) -> bool:
    """Semantic validation turns malformed local output into an in-process repair attempt."""

    value = payload.get(field)
    return isinstance(value, str) and bool(value.strip())


def _has_nonempty_list(payload: dict[str, Any], field: str) -> bool:
    """Required summary sections cannot silently collapse to empty arrays."""

    value = payload.get(field)
    return isinstance(value, list) and bool(value)


def _dedupe_text_items(value: Any) -> list[str]:
    """Keep model-generated summary sections readable when a provider repeats a point."""

    if not isinstance(value, list):
        return []
    result: list[str] = []
    seen: set[str] = set()
    for item in value:
        if not isinstance(item, str):
            continue
        text = item.strip()
        key = " ".join(text.split()).casefold()
        if not text or key in seen:
            continue
        seen.add(key)
        result.append(text)
    return result


def _has_explanation_shape(payload: dict[str, Any]) -> bool:
    """Require a real summary while allowing trusted normalization of omitted optional sections."""

    if not _has_nonempty_string(payload, "summary"):
        return False
    for field in (
        "missing_context",
        "rare_terms",
        "possible_asr_errors",
        "review_questions",
    ):
        if field in payload and not isinstance(payload[field], list):
            return False
    return "confidence" not in payload or isinstance(payload["confidence"], int | float)


def _uses_requested_explanation_language(
    payload: dict[str, Any], target_language: str
) -> bool:
    """A Chinese session never accepts an English-only summary as a completed card."""

    if not target_language.lower().startswith("zh"):
        return True
    summary = payload.get("summary")
    return isinstance(summary, str) and any("\u4e00" <= char <= "\u9fff" for char in summary)


def _normalize_explanation(
    raw: dict[str, Any], segment_ids: list[str], page_ids: list[str]
) -> ExplanationResponse | None:
    """Pydantic validates structure while local allowlists remove fabricated evidence IDs."""

    normalized = {
        "summary": raw.get("summary", ""),
        "missing_context": raw.get("missing_context", []),
        "rare_terms": raw.get("rare_terms", []),
        "possible_asr_errors": raw.get("possible_asr_errors", []),
        "review_questions": raw.get("review_questions", []),
        "evidence_segment_ids": raw.get("evidence_segment_ids", segment_ids),
        "asset_page_ids": raw.get("asset_page_ids", page_ids),
        "confidence": raw.get("confidence", 0.5),
        "provider": "pending",
    }
    try:
        candidate = ExplanationResponse.model_validate(normalized)
    except ValueError:
        return None
    allowed_segments = set(segment_ids)
    allowed_pages = set(page_ids)
    candidate.evidence_segment_ids = [
        item for item in candidate.evidence_segment_ids if item in allowed_segments
    ]
    candidate.asset_page_ids = [item for item in candidate.asset_page_ids if item in allowed_pages]
    for context_item in candidate.missing_context:
        context_item.evidence_segment_ids = [
            value for value in context_item.evidence_segment_ids if value in allowed_segments
        ]
    for term_item in candidate.rare_terms:
        term_item.evidence_segment_ids = [
            value for value in term_item.evidence_segment_ids if value in allowed_segments
        ]
        term_item.asset_page_ids = [
            value for value in term_item.asset_page_ids if value in allowed_pages
        ]
    if not candidate.evidence_segment_ids:
        candidate.evidence_segment_ids = segment_ids
    return candidate


def _parse_asset_sync(suffix: str, data: bytes) -> AssetParseResponse:
    """Format-specific parsers extract page text before any optional OCR or VLM work."""

    if suffix == ".pptx":
        presentation = Presentation(io.BytesIO(data))
        pages = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text")]
            nonempty = [text for text in texts if text]
            pages.append(
                ParsedPage(
                    page_number=index,
                    title=nonempty[0] if nonempty else f"Slide {index}",
                    text="\n".join(nonempty),
                )
            )
        return AssetParseResponse(parser="python-pptx:1", pages=pages)
    if suffix == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        pages = [
            ParsedPage(
                page_number=index,
                title=f"Page {index}",
                text=(page.extract_text() or "").strip(),
            )
            for index, page in enumerate(reader.pages, start=1)
        ]
        return AssetParseResponse(parser="pypdf:1", pages=pages)
    if suffix == ".docx":
        document = Document(io.BytesIO(data))
        text = "\n".join(
            paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text
        )
        return AssetParseResponse(
            parser="python-docx:1",
            pages=[ParsedPage(page_number=1, title="Document", text=text)],
        )
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        with Image.open(io.BytesIO(data)) as image:
            title = f"Image {image.width}×{image.height}"
        return AssetParseResponse(
            parser="pillow-metadata:1",
            pages=[ParsedPage(page_number=1, title=title, text="")],
        )
    if suffix in {".txt", ".md", ".csv"}:
        text = data.decode("utf-8", errors="replace")
        return AssetParseResponse(
            parser="utf8-text:1",
            pages=[ParsedPage(page_number=1, title="Text", text=text)],
        )
    raise ValueError(f"unsupported asset suffix: {suffix}")


if __name__ == "__main__":
    # Direct execution supports local work; packaged launches reuse the same app object.
    import uvicorn

    uvicorn.run(app, host="127.0.0.1", port=8790)
